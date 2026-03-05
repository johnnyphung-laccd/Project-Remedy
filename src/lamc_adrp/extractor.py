"""Stage 3: Content Extraction — OCR and native parsing for all document types.

Routes each document to the appropriate extraction strategy:
  - PDF  -> GLM-OCR layout parsing
  - DOCX -> python-docx native parser (GLM-OCR fallback for image-heavy docs)
  - PPTX -> python-pptx native parser (GLM-OCR fallback)
  - XLSX -> openpyxl native parser (sheet-by-sheet markdown tables)

All extractors emit clean Markdown preserving headings, lists, tables,
emphasis, and image placeholders.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from lamc_adrp.config import PipelineConfig
from lamc_adrp.database import DatabaseManager
from lamc_adrp.models import DocumentJob, FileType, JobStatus
from lamc_adrp.zai_client import ZAIClient, ZAIClientError

logger = logging.getLogger(__name__)

# Threshold: if a DOCX/PPTX has this many embedded images relative to
# text paragraphs, we fall back to GLM-OCR for more accurate extraction.
_IMAGE_HEAVY_RATIO = 0.5


class ExtractionError(Exception):
    """Raised when content extraction fails for a document."""


class ContentExtractor:
    """Orchestrates content extraction across all supported file types.

    Parameters
    ----------
    config:
        Pipeline configuration.
    zai:
        Initialised Z.AI API client.
    db:
        Database manager for persisting job state.
    """

    def __init__(
        self,
        config: PipelineConfig,
        zai: ZAIClient,
        db: DatabaseManager,
    ) -> None:
        self._config = config
        self._zai = zai
        self._db = db

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------

    async def extract(self, job: DocumentJob) -> str:
        """Extract content from the downloaded document.

        Updates the job status through EXTRACTING -> EXTRACTED (or FAILED)
        and stores the result in ``job.ocr_markdown``.

        Returns
        -------
        str
            Extracted Markdown content.
        """
        job.status = JobStatus.EXTRACTING
        await self._db.update_job(job)
        logger.info("Extracting content for job %s (%s)", job.id, job.file_type)

        try:
            file_path = Path(job.local_path)
            if not file_path.exists():
                raise ExtractionError(f"Local file not found: {file_path}")

            file_type = job.file_type
            if file_type is None:
                raise ExtractionError("Job has no file_type set.")

            if file_type == FileType.PDF:
                markdown = await self._extract_pdf(file_path)
            elif file_type in (FileType.DOCX, FileType.DOC):
                markdown = await self._extract_docx(file_path)
            elif file_type in (FileType.PPTX, FileType.PPT):
                markdown = await self._extract_pptx(file_path)
            elif file_type in (FileType.XLSX, FileType.XLS):
                markdown = await self._extract_xlsx(file_path)
            else:
                raise ExtractionError(f"Unsupported file type: {file_type}")

            if not markdown.strip():
                raise ExtractionError("Extraction produced empty content.")

            job.ocr_markdown = markdown
            job.status = JobStatus.EXTRACTED
            await self._db.update_job(job)
            logger.info(
                "Extraction complete for job %s — %d chars of markdown",
                job.id,
                len(markdown),
            )
            return markdown

        except Exception as exc:
            error_msg = f"Extraction failed: {exc}"
            logger.error("Job %s: %s", job.id, error_msg)
            job.status = JobStatus.FAILED
            job.error_message = error_msg
            await self._db.update_job(job)
            raise ExtractionError(error_msg) from exc

    # ------------------------------------------------------------------
    # PDF extraction
    # ------------------------------------------------------------------

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract PDF content via GLM-OCR."""
        logger.debug("PDF extraction via GLM-OCR: %s", file_path.name)
        return await self._zai.ocr(file_path=file_path)

    # ------------------------------------------------------------------
    # DOCX extraction
    # ------------------------------------------------------------------

    async def _extract_docx(self, file_path: Path) -> str:
        """Extract DOCX content using python-docx, with GLM-OCR fallback."""
        try:
            from docx import Document as DocxDocument  # type: ignore[import-untyped]
            from docx.opc.constants import RELATIONSHIP_TYPE as RT  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "python-docx not installed — falling back to GLM-OCR for %s",
                file_path.name,
            )
            return await self._zai.ocr(file_path=file_path)

        logger.debug("DOCX native extraction: %s", file_path.name)
        doc = DocxDocument(str(file_path))

        # Count images to decide if we should fall back to OCR.
        image_count = sum(
            1
            for rel in doc.part.rels.values()
            if "image" in rel.reltype
        )
        paragraph_count = len(doc.paragraphs)
        if paragraph_count > 0 and image_count / paragraph_count > _IMAGE_HEAVY_RATIO:
            logger.info(
                "DOCX %s is image-heavy (%d images / %d paragraphs) "
                "— falling back to GLM-OCR",
                file_path.name,
                image_count,
                paragraph_count,
            )
            return await self._zai.ocr(file_path=file_path)

        parts: list[str] = []
        image_index = 0

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                text = self._docx_paragraph_to_md(element, doc)
                if text:
                    parts.append(text)

            elif tag == "tbl":
                table_md = self._docx_table_to_md(element, doc)
                if table_md:
                    parts.append(table_md)

        # Append image placeholders.
        for i in range(image_count):
            parts.append(f"\n![Image {i + 1}](image_{i + 1}.png)\n")

        return "\n\n".join(parts)

    def _docx_paragraph_to_md(self, para_element: Any, doc: Any) -> str:
        """Convert a DOCX paragraph XML element to Markdown."""
        from docx.text.paragraph import Paragraph  # type: ignore[import-untyped]

        para = Paragraph(para_element, doc)
        text = ""

        # Gather run-level formatting.
        for run in para.runs:
            run_text = run.text or ""
            if not run_text:
                continue
            if run.bold:
                run_text = f"**{run_text}**"
            if run.italic:
                run_text = f"*{run_text}*"
            text += run_text

        text = text.strip()
        if not text:
            return ""

        # Determine heading level from style name.
        style_name = (para.style.name or "").lower() if para.style else ""
        if style_name.startswith("heading"):
            # Extract heading level number.
            match = re.search(r"\d+", style_name)
            level = int(match.group()) if match else 1
            level = min(level, 6)
            return f"{'#' * level} {text}"

        # Detect list styles.
        if style_name.startswith("list bullet") or style_name.startswith("list"):
            return f"- {text}"
        if style_name.startswith("list number"):
            return f"1. {text}"

        return text

    def _docx_table_to_md(self, tbl_element: Any, doc: Any) -> str:
        """Convert a DOCX table XML element to a Markdown table."""
        from docx.table import Table  # type: ignore[import-untyped]

        table = Table(tbl_element, doc)
        rows = table.rows
        if not rows:
            return ""

        md_rows: list[str] = []
        for i, row in enumerate(rows):
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            md_rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

        return "\n".join(md_rows)

    # ------------------------------------------------------------------
    # PPTX extraction
    # ------------------------------------------------------------------

    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract PPTX content slide-by-slide using python-pptx."""
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "python-pptx not installed — falling back to GLM-OCR for %s",
                file_path.name,
            )
            return await self._zai.ocr(file_path=file_path)

        logger.debug("PPTX native extraction: %s", file_path.name)
        prs = Presentation(str(file_path))

        # Check if image-heavy.
        total_shapes = 0
        image_shapes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                if shape.shape_type and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_shapes += 1

        if total_shapes > 0 and image_shapes / total_shapes > _IMAGE_HEAVY_RATIO:
            logger.info(
                "PPTX %s is image-heavy (%d/%d shapes are images) "
                "— falling back to GLM-OCR",
                file_path.name,
                image_shapes,
                total_shapes,
            )
            return await self._zai.ocr(file_path=file_path)

        parts: list[str] = []
        image_index = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"## Slide {slide_num}"]

            for shape in slide.shapes:
                # Title shape.
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = self._pptx_paragraph_to_md(para, shape)
                        if text:
                            slide_parts.append(text)

                # Table shape.
                if shape.has_table:
                    table_md = self._pptx_table_to_md(shape.table)
                    if table_md:
                        slide_parts.append(table_md)

                # Image shape.
                if shape.shape_type and shape.shape_type == 13:
                    image_index += 1
                    slide_parts.append(
                        f"![Slide {slide_num} Image {image_index}]"
                        f"(slide_{slide_num}_image_{image_index}.png)"
                    )

            # Slide notes.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Speaker Notes:** {notes_text}")

            parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(parts)

    def _pptx_paragraph_to_md(self, para: Any, shape: Any) -> str:
        """Convert a python-pptx paragraph to Markdown."""
        text = ""
        for run in para.runs:
            run_text = run.text or ""
            if not run_text:
                continue
            if run.font.bold:
                run_text = f"**{run_text}**"
            if run.font.italic:
                run_text = f"*{run_text}*"
            text += run_text

        text = text.strip()
        if not text:
            return ""

        # Check if this is the title placeholder.
        is_title = False
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
            if ph_idx in (0, 1):  # Title or Center Title
                is_title = True

        if is_title:
            return f"### {text}"

        # Bullet level.
        level = para.level if hasattr(para, "level") and para.level else 0
        if level > 0:
            indent = "  " * (level - 1)
            return f"{indent}- {text}"

        return text

    def _pptx_table_to_md(self, table: Any) -> str:
        """Convert a python-pptx table to a Markdown table."""
        rows = table.rows
        if not rows:
            return ""

        md_rows: list[str] = []
        for i, row in enumerate(rows):
            cells = [
                cell.text.strip().replace("|", "\\|") for cell in row.cells
            ]
            md_rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

        return "\n".join(md_rows)

    # ------------------------------------------------------------------
    # XLSX extraction
    # ------------------------------------------------------------------

    async def _extract_xlsx(self, file_path: Path) -> str:
        """Extract XLSX content sheet-by-sheet as Markdown tables."""
        try:
            from openpyxl import load_workbook  # type: ignore[import-untyped]
        except ImportError:
            logger.warning(
                "openpyxl not installed — falling back to GLM-OCR for %s",
                file_path.name,
            )
            return await self._zai.ocr(file_path=file_path)

        logger.debug("XLSX native extraction: %s", file_path.name)
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts: list[str] = [f"## {sheet_name}"]

            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                cells = [
                    str(cell).strip() if cell is not None else ""
                    for cell in row
                ]
                # Skip completely empty rows.
                if any(cells):
                    rows_data.append(cells)

            if not rows_data:
                sheet_parts.append("*Empty sheet*")
                parts.append("\n\n".join(sheet_parts))
                continue

            # Normalise column count (pad shorter rows).
            max_cols = max(len(r) for r in rows_data)
            for row in rows_data:
                while len(row) < max_cols:
                    row.append("")

            # Build Markdown table.
            for i, row in enumerate(rows_data):
                escaped = [c.replace("|", "\\|") for c in row]
                sheet_parts.append("| " + " | ".join(escaped) + " |")
                if i == 0:
                    sheet_parts.append(
                        "| " + " | ".join("---" for _ in escaped) + " |"
                    )

            parts.append("\n".join(sheet_parts))

        wb.close()
        return "\n\n---\n\n".join(parts)
